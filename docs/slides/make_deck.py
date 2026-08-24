"""Results deck for the WACV thermal-deer paper.

Every number here is copied from overleaf_WACV/, so the talk and the manuscript
cannot drift apart. Nothing is computed fresh.
"""
import os
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIG = "/work/nvme/bgte/tislam6/wildlife_project/overleaf_WACV/figures"
AST = "/work/nvme/bgte/tislam6/wildlife_project/docs/slides/assets"
OUT = "/work/nvme/bgte/tislam6/wildlife_project/docs/slides/wildlife_results.pptx"

PAPER  = C(0xFA, 0xF8, 0xF7)
INK    = C(0x24, 0x1C, 0x1E)
MUTE   = C(0x7A, 0x6E, 0x70)
RULE   = C(0xE4, 0xDD, 0xDE)
MAROON = C(0x8C, 0x2D, 0x3F)
DEEP   = C(0x4A, 0x1B, 0x33)
EMBER  = C(0xC4, 0x45, 0x3C)

prs = Presentation()
prs.slide_width, prs.slide_height = I(13.333), I(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5
M = 0.72                       # page margin


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s


def text(s, x, y, w, h, body, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         spacing=1.0, font="Calibri", anchor=MSO_ANCHOR.TOP, gap=0):
    tb = s.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(body if isinstance(body, list) else [body]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        if i and gap:
            p.space_before = Pt(gap)
        txt, sz, cl, bd = (line if isinstance(line, tuple) else (line, size, color, bold))
        run = p.add_run(); run.text = txt
        run.font.size = Pt(sz); run.font.color.rgb = cl
        run.font.bold = bd; run.font.name = font
    return tb


def bar(s, x, y, w, h, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def header(s, kicker, title):
    bar(s, M, 0.62, 0.055, 0.52, MAROON)
    text(s, M + 0.22, 0.60, 9.0, 0.3, kicker.upper(), size=11.5, color=MAROON, bold=True)
    text(s, M + 0.22, 0.86, 11.6, 0.5, title, size=25, bold=True)


def picture_fit(s, path, x, y, w, h):
    """Contain-fit, so nothing is ever stretched or cropped."""
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return s.shapes.add_picture(path, I(x + (w - pw) / 2), I(y + (h - ph) / 2),
                                I(pw), I(ph))


def kpi(s, x, y, w, value, label, color=MAROON):
    bar(s, x, y, 0.045, 1.28, color)
    text(s, x + 0.19, y - 0.06, w, 0.62, value, size=38, bold=True, color=INK)
    text(s, x + 0.19, y + 0.62, w, 0.66, label, size=12.5, color=MUTE, spacing=1.15)


def table(s, x, y, w, cols, rows, widths, size=12, head_size=11, row_h=0.285,
          bold_rows=(), right_from=1, head_gap=0.055):
    """Flat table: a maroon rule under the header, hairlines between rows, no boxes.
    `widths` are relative and are normalised to `w`."""
    tot = float(sum(widths))
    xs, acc = [], 0.0
    for cw in widths:
        xs.append(x + acc / tot * w); acc += cw
    ws = [cw / tot * w for cw in widths]

    for j, (cx, cw, c) in enumerate(zip(xs, ws, cols)):
        al = PP_ALIGN.LEFT if j < right_from else PP_ALIGN.RIGHT
        text(s, cx, y, cw, 0.26, c, size=head_size, bold=True, color=MUTE, align=al)
    bar(s, x, y + 0.27, w, 0.018, MAROON)

    yy = y + 0.27 + head_gap
    for i, r in enumerate(rows):
        strong = i in bold_rows
        for j, (cx, cw, v) in enumerate(zip(xs, ws, r)):
            al = PP_ALIGN.LEFT if j < right_from else PP_ALIGN.RIGHT
            text(s, cx, yy + 0.045, cw, row_h, str(v), size=size,
                 bold=strong, color=INK if strong else INK, align=al)
        yy += row_h
        if i < len(rows) - 1 and (i + 1) not in bold_rows:
            bar(s, x, yy, w, 0.007, RULE)
        elif i < len(rows) - 1:
            bar(s, x, yy, w, 0.014, RULE)
    return yy


# ══════════════════════════════════════════════════ 1 — title
s = slide()
bar(s, 0, 0, W, 0.11, MAROON)
text(s, M, 2.05, 11.4, 0.34, "WACV 2027 SUBMISSION  ·  RESULTS TO DATE",
     size=13, color=MAROON, bold=True)
text(s, M, 2.55, 11.6, 1.9,
     "Counting Is Not Detecting", size=52, bold=True)
text(s, M, 3.62, 11.6, 0.9,
     "Localizing the bottleneck in thermal wildlife video", size=25, color=MUTE)
bar(s, M, 4.72, 3.0, 0.02, RULE)
text(s, M, 5.00, 11.6, 1.2,
     ["White-tailed deer detection and counting from vehicle-mounted FLIR thermal transects",
      ("32 videos  ·  521,930 frames  ·  236 individually tracked animals  ·  4 sites, southern Illinois",
       15, MUTE, False)],
     size=15, color=INK, spacing=1.35, gap=9)

# ══════════════════════════════════════════════════ 2 — pipeline
s = slide()
header(s, "Approach", "One detector, one tracker, a three-parameter confirmation rule")
picture_fit(s, f"{FIG}/wildlife_pipeline.png", M, 1.95, W - 2 * M, 2.5)
bar(s, M, 4.78, W - 2 * M, 0.02, RULE)
cw = (W - 2 * M - 1.0) / 3
for i, (t, b) in enumerate([
    ("Detection", "YOLO11m @ 640 px, chosen from a 12-architecture benchmark on the "
                  "criterion that matters downstream, not on mAP."),
    ("Association", "BoT-SORT with global motion compensation — mandatory, the vehicle "
                    "travels 25–90 km/h and the whole frame translates."),
    ("Confirmation", "Accept a track on length, span and top-5 mean confidence. Three "
                     "parameters, swept on training video and frozen."),
]):
    x = M + i * (cw + 0.5)
    text(s, x, 5.05, cw, 0.32, t, size=15.5, bold=True, color=MAROON)
    text(s, x, 5.44, cw, 1.5, b, size=12.5, color=INK, spacing=1.3)

# ══════════════════════════════════════════════════ 3 — dataset
s = slide()
header(s, "Dataset", "32 annotated thermal transects, 236 individually tracked animals")
picture_fit(s, f"{AST}/dataset_grid.png", M, 1.80, W - 2 * M, 2.42)
text(s, M, 4.32, W - 2 * M, 0.3,
     "Top: frames as the detector receives them, after CLAHE normalization.   "
     "Bottom: the same frames with their annotation.",
     size=11.5, color=MUTE)

table(s, M, 4.86, 6.35,
      ["Split", "Videos", "Frames", "With deer", "Empty", "Boxes"],
      [["Train", "19", "12,852", "5,013", "7,839", "7,547"],
       ["Val",   "4",  "2,337",  "935",   "1,402", "2,564"],
       ["Test",  "9",  "3,725",  "1,490", "2,235", "2,792"],
       ["Total", "32", "18,914", "7,438", "11,476", "12,903"]],
      widths=[1.5, 1.0, 1.15, 1.25, 1.0, 1.1],
      size=12.5, bold_rows=(3,))

text(s, 7.55, 4.86, 5.05, 0.3, "What matters here", size=14.5, bold=True, color=MAROON)
text(s, 7.55, 5.28, 5.05, 1.9,
     ["11,476 frames are explicit negatives. A survey system must not fire on warm rocks "
      "or structures, so background frames are annotated, not discarded.",
      "Splits are by video and site-stratified — no frame from a training video appears "
      "in evaluation."],
     size=13, spacing=1.28, gap=9)

# ══════════════════════════════════════════════════ 4 — detection benchmark
s = slide()
header(s, "Detection", "Twelve architectures, and very little separates them")
DET = [["YOLOv8m",  "0.459", "0.780", "0.572", "0.660", "0.727"],
       ["ATSS R50", "0.380", "0.683", "0.635", "0.658", "0.673"],
       ["TOOD R50", "0.396", "0.683", "0.621", "0.650", "0.670"],
       ["YOLO11m",  "0.460", "0.897", "0.503", "0.645", "0.775"],
       ["RT-DETR-L","0.434", "0.643", "0.605", "0.623", "0.635"],
       ["YOLO12m",  "0.508", "0.929", "0.452", "0.608", "0.767"],
       ["YOLOv9m",  "0.506", "0.917", "0.450", "0.604", "0.759"],
       ["YOLOv10m", "0.510", "0.874", "0.453", "0.597", "0.737"],
       ["DINO R50", "0.365", "0.560", "0.572", "0.566", "0.562"],
       ["Faster R-CNN", "0.418", "0.393", "0.723", "0.509", "0.433"],
       ["RTMDet-m", "0.466", "0.325", "0.779", "0.459", "0.368"]]
table(s, M, 1.92, 6.0,
      ["Model", "AP50", "Precision", "Recall", "F1", "F0.5"],
      DET, widths=[1.75, 0.85, 1.15, 0.9, 0.8, 0.85],
      size=11.5, head_size=10.5, row_h=0.375, bold_rows=(3,))

picture_fit(s, f"{AST}/detect_grid.png", 7.15, 1.92, 5.45, 2.75)
text(s, 7.15, 4.80, 5.45, 0.3, "Why YOLO11m", size=14.5, bold=True, color=MAROON)
text(s, 7.15, 5.22, 5.45, 1.9,
     ["AP50 spans only 0.365–0.510 across designs as different as Faster R-CNN, DINO and "
      "YOLO12m — the signature of a data-limited rather than an architecture-limited regime.",
      "The confirmation stage downstream must reject tens of thousands of false candidates "
      "against a few hundred true ones, so we select on precision-weighted F0.5, where "
      "YOLO11m ranks first at 0.775."],
     size=12.5, spacing=1.26, gap=9)

# ══════════════════════════════════════════════════ 5 — headline numbers
s = slide()
header(s, "Results", "Detection is not the bottleneck")
for i, (v, l, c) in enumerate([
    ("98.8%", "of held-out animals\nare detected at all", DEEP),
    ("88.0%", "acquire their own\ndistinct track", MAROON),
    ("62.7%", "survive confirmation\nand are counted", EMBER),
    ("2.38", "mean absolute error,\nanimals per video", MUTE),
]):
    kpi(s, M + i * 3.06, 2.05, 2.7, v, l, c)
bar(s, M, 3.66, W - 2 * M, 0.02, RULE)
picture_fit(s, f"{AST}/funnel.png", M - 0.1, 3.95, 6.4, 3.05)
text(s, 7.15, 3.95, 5.45, 0.34, "What this says", size=15.5, bold=True, color=MAROON)
text(s, 7.15, 4.38, 5.45, 2.2,
     ["Nearly every animal is found. The loss is concentrated after detection — "
      "in association, and above all in the decision of which candidate tracks are real.",
      "A better backbone cannot recover the 36 points between detected and counted, "
      "because those animals were never missed by the detector in the first place."],
     size=13.5, spacing=1.28, gap=11)
text(s, 7.15, 6.72, 5.45, 0.3,
     "Held out: rule swept on 19 training videos, frozen, reported on 13 unseen.",
     size=11, color=MUTE)

# ══════════════════════════════════════════════════ 6 — per-video breakdown
s = slide()
header(s, "Per-video breakdown", "Reading a row left to right attributes that video's loss to a stage")
PV = [["NShelbyRd (blue)", "SHB", "27", "22", "16", "20", "−7"],
      ["GolfDr",           "SHB", "12", "12", "10", "14", "+2"],
      ["NWolfCreek (orange)","SHB","9",  "9",  "9",  "5",  "−4"],
      ["GiantCityRd",      "TON", "8",  "7",  "7",  "3",  "−5"],
      ["Robinson",         "SHW", "7",  "7",  "7",  "7",  "0"],
      ["Melvin",           "SHW", "5",  "5",  "5",  "1",  "−4"],
      ["AquaCultureRd",    "TON", "3",  "1",  "1",  "0",  "−3"],
      ["N25thBlue",        "MAS", "3",  "3",  "3",  "3",  "0"],
      ["NMarseilles",      "MAS", "3",  "3",  "3",  "1",  "−2"],
      ["OikosRd",          "TON", "2",  "2",  "2",  "1",  "−1"],
      ["TouchofNature",    "TON", "2",  "1",  "1",  "0",  "−2"],
      ["ChipsRd",          "TON", "1",  "1",  "1",  "1",  "0"],
      ["SIron",            "SHW", "1",  "1",  "1",  "2",  "+1"],
      ["Total",            "",    "83", "74", "66", "58", "−25"]]
table(s, M, 1.88, 7.15,
      ["Video", "Site", "Animals", "Detected", "Own track", "Counted", "Error"],
      PV, widths=[2.35, 0.75, 1.05, 1.05, 1.1, 0.95, 0.8],
      size=11.5, head_size=10.5, row_h=0.335, right_from=2, bold_rows=(13,))

text(s, 8.55, 1.88, 4.05, 0.3, "How to read it", size=14.5, bold=True, color=MAROON)
text(s, 8.55, 2.30, 4.05, 4.4,
     ["Animals is the annotated count. Detected and Own track are the two intermediate "
      "stages; Counted is what the system reports.",
      "The error concentrates in dense-group footage. NShelbyRd loses 5 animals to "
      "detection and a further 6 to association — the largest tracking failure in the "
      "corpus, and the reason group video dominates the total.",
      "Only two videos over-count, and both do so with their detection and association "
      "columns complete. Their error is confirmation accepting spurious tracks, not "
      "anything upstream.",
      "MAE 2.38 animals per video; the system under-counts and essentially never inflates."],
     size=12.5, spacing=1.26, gap=10)

# ══════════════════════════════════════════════════ 7 — qualitative
s = slide()
header(s, "Qualitative results", "Tracking holds through overlap and low contrast")
picture_fit(s, f"{FIG}/qualitative.png", M, 1.85, W - 2 * M, 3.55)
bar(s, M, 5.62, W - 2 * M, 0.02, RULE)
cw = (W - 2 * M - 1.0) / 3
for i, (t, b) in enumerate([
    ("Identities persist", "Top row is the same scene 0.3 s apart. Tracks keep their "
                           "identities while the camera translates beneath them."),
    ("Overlap is resolved", "Animals whose boxes intersect are still carried as separate "
                            "tracks — the failure that silently merges two deer into one."),
    ("Confidence is honest", "The per-track score falls with range, from a near animal at "
                             "87 px to a distant one at 22 px in a single frame."),
]):
    x = M + i * (cw + 0.5)
    text(s, x, 5.86, cw, 0.3, t, size=14.5, bold=True, color=MAROON)
    text(s, x, 6.22, cw, 1.0, b, size=12, color=INK, spacing=1.28)

# ══════════════════════════════════════════════════ 8 — the inversion
s = slide()
header(s, "Central finding", "Improving candidate generation makes counting worse")
picture_fit(s, f"{AST}/inversion.png", M - 0.1, 2.0, 6.5, 3.2)
text(s, M + 0.15, 5.42, 6.2, 0.9,
     "Four independent interventions each enlarged the candidate pool as designed. "
     "Every one of them counted fewer animals.",
     size=13, color=MUTE, spacing=1.3)
text(s, 7.35, 2.05, 5.25, 0.34, "Why", size=15.5, bold=True, color=MAROON)
text(s, 7.35, 2.48, 5.25, 2.2,
     "Each intervention adds far more false candidates than real ones — the richest pool "
     "holds 219 true tracks against 27,460 false. A rule tuned to survive that ratio has "
     "to be aggressive enough to reject real animals along with the noise.",
     size=14, spacing=1.32)
bar(s, 7.35, 4.72, 5.25, 0.02, RULE)
text(s, 7.35, 4.98, 5.25, 0.34, "And it resists learning", size=15.5, bold=True, color=MAROON)
text(s, 7.35, 5.41, 5.25, 1.8,
     "A transformer, gradient boosting and logistic regression all lose to the "
     "three-parameter rule under three separate protocols. Error grows with model "
     "capacity — the constraint is supervision, not architecture.",
     size=14, spacing=1.32)

# ══════════════════════════════════════════════════ 9 — takeaways
s = slide()
header(s, "Summary", "What we can claim today")
rows = [
    ("Corpus", "32 fully annotated thermal transects, 236 individually tracked deer, "
               "21,646 boxes. Released with the paper."),
    ("Detector", "12 architectures benchmarked under four matching criteria. YOLO11m @ 640 px, "
                 "0.897 precision, selected on the criterion the counter actually needs."),
    ("Preprocessing", "CLAHE normalization is worth more than any architectural choice we "
                      "tested — test mAP50 0.519 with it, 0.299 without."),
    ("Counting", "98.8% detected, 88.0% distinctly tracked, 62.7% counted on video the "
                 "detector never saw. MAE 2.38 animals per video, and the system under-counts "
                 "rather than over-counts."),
    ("Auditability", "Every counted animal carries a calibrated score and a contact sheet of "
                     "its own frames; 92.1% score above 0.80."),
]
y = 1.95
for t, b in rows:
    bar(s, M, y + 0.06, 0.035, 0.46, MAROON)
    text(s, M + 0.2, y, 2.15, 0.4, t, size=14.5, bold=True, color=MAROON)
    text(s, M + 2.5, y - 0.02, W - 2 * M - 2.5, 0.9, b, size=13.5, spacing=1.28)
    y += 0.94
bar(s, M, 6.72, W - 2 * M, 0.02, RULE)
text(s, M, 6.92, W - 2 * M, 0.34,
     "Manuscript drafted end to end; remaining work is length and final ablations.",
     size=12.5, color=MUTE)

# ══════════════════════════════════════════════════ 10 — roadmap
s = slide()
header(s, "What comes next", "Four workstreams")
items = [
    ("1", "Counting paper", "Thermal transect detection and counting. Drafted end to end; "
                            "remaining work is length and final ablations.", "In progress", EMBER),
    ("2", "Vision–language scene understanding",
          "Move from a number to a description: how many, how far, what are they doing, "
          "what is occluding them.", "Proposed", MAROON),
    ("3", "Edge deployment",
          "Run the detector on the camera rather than on a cluster, so a trail camera or "
          "drone reports counts in the field.", "Proposed", MAROON),
    ("4", "Collaborative grant writing",
          "Convert the corpus, the pipeline and the evaluation framework into funded, "
          "multi-institution work.", "Proposed", DEEP),
]
y = 1.95
for n, t, b, status, col in items:
    bar(s, M, y, 0.05, 1.02, col)
    text(s, M + 0.26, y - 0.04, 0.5, 0.45, n, size=27, bold=True, color=col)
    text(s, M + 0.95, y - 0.02, 7.4, 0.34, t, size=17, bold=True)
    text(s, M + 0.95, y + 0.36, 7.9, 0.66, b, size=13, color=MUTE, spacing=1.26)
    text(s, 10.9, y + 0.02, 1.7, 0.3, status, size=12.5, bold=True, color=col,
         align=PP_ALIGN.RIGHT)
    y += 1.22

# ══════════════════════════════════════════════════ 11 — VLM
s = slide()
header(s, "Workstream 2 · proposed", "From a count to an explanation of the scene")
picture_fit(s, f"{AST}/vlm_scene.png", M, 1.82, 7.2, 4.15)
text(s, M, 6.06, 7.2, 0.62,
     "Trail-camera frame, TON11. The same scene our thermal pipeline would reduce to the "
     "single number 4. Boxes and ranges here illustrate the proposed output.",
     size=11.5, color=MUTE, spacing=1.25)

X = 8.25
text(s, X, 1.82, 4.35, 0.3, "Proposed model output", size=15, bold=True, color=MAROON)
text(s, X, 2.24, 4.35, 2.5,
     ["\u201cFour white-tailed deer. One adult foraging in the immediate foreground, two "
      "browsing together at mid-range and partly occluded by saplings, one adult at the "
      "far right beside a mature oak.",
      "All four are head-down feeding, none alert, no antlers visible. The group is "
      "dispersed and stationary rather than travelling.\u201d"],
     size=13, spacing=1.3, gap=8)
bar(s, X, 4.62, 4.35, 0.018, RULE)
text(s, X, 4.82, 4.35, 0.3, "Structured fields alongside the prose", size=13.5, bold=True,
     color=MAROON)
for i, (k, v) in enumerate([("count", "4"), ("range per animal", "~4, 14, 16, 19 m"),
                            ("mask", "per-animal segmentation"),
                            ("behaviour", "foraging, not alert"),
                            ("occlusion", "2 of 4 partly hidden")]):
    yy = 5.22 + i * 0.30
    text(s, X, yy, 1.85, 0.28, k, size=12, color=MUTE)
    text(s, X + 1.9, yy, 2.45, 0.28, v, size=12, bold=True)
text(s, X, 6.82, 4.35, 0.3,
     "Illustrative — a target capability, not a current result.",
     size=11, color=MUTE)

# ══════════════════════════════════════════════════ 12 — edge
s = slide()
header(s, "Workstream 3 · proposed", "Move the detector onto the camera")
text(s, M, 1.92, 5.6, 0.3, "Why it is plausible now", size=15, bold=True, color=MAROON)
text(s, M, 2.34, 5.6, 2.6,
     ["The production detector is a 40\u2009MB checkpoint running at 640\u2009px on a "
      "single class. That is a small model by current standards, and the pipeline around "
      "it — motion compensation, association, the confirmation rule — carries no learned "
      "weights and negligible compute.",
      "Nothing in the design assumes a cluster. The counting run is embarrassingly "
      "parallel across videos because each video is independent."],
     size=13.5, spacing=1.3, gap=10)
bar(s, M, 5.05, 5.6, 0.018, RULE)
text(s, M, 5.25, 5.6, 0.3, "What changes for the survey", size=15, bold=True, color=MAROON)
text(s, M, 5.67, 5.6, 1.4,
     "A camera that reports a count instead of storing footage removes the retrieval and "
     "manual-review step that currently bounds how many sites a survey can cover.",
     size=13.5, spacing=1.3)

text(s, 7.05, 1.92, 5.55, 0.3, "Plan", size=15, bold=True, color=MAROON)
steps = [("Export", "TensorRT / ONNX, INT8 and FP16 variants"),
         ("Measure", "throughput, latency and power on the lab Jetson — "
                     "no on-device numbers exist yet"),
         ("Verify", "confirm counting accuracy is unchanged after quantization, "
                   "using the same held-out protocol"),
         ("Field", "battery-powered trail-camera unit reporting counts over LTE")]
y = 2.34
for i, (t, b) in enumerate(steps):
    bar(s, 7.05, y + 0.04, 0.035, 0.42, MAROON)
    text(s, 7.25, y, 1.28, 0.3, t, size=13.5, bold=True, color=MAROON)
    text(s, 8.6, y - 0.02, 4.0, 0.9, b, size=12.5, spacing=1.26)
    y += 0.98
text(s, 7.05, 6.5, 5.55, 0.5,
     "Quantization is the open question: the animals are already at the sensor's "
     "resolution limit, so reduced precision may cost recall where it is scarcest.",
     size=12, color=MUTE, spacing=1.28)

# ══════════════════════════════════════════════════ 13 — collaboration
s = slide()
header(s, "Workstream 4 · proposed", "What we can bring to a collaborative proposal")
assets = [
    ("A released corpus", "32 annotated thermal transects with per-individual tracks — "
                          "rare enough that the scarcity is itself a contribution."),
    ("An evaluation framework", "The reach / own-track / counted decomposition transfers "
                                "to any survey modality, not just thermal."),
    ("A negative result worth citing", "Improving candidate generation degrades counting. "
                                       "That shapes how a funded effort should spend its "
                                       "money."),
    ("A deployment path", "Workstreams 2 and 3 turn a measurement into an instrument "
                          "ecologists can actually operate."),
]
y = 1.95
for t, b in assets:
    bar(s, M, y + 0.05, 0.035, 0.5, MAROON)
    text(s, M + 0.2, y, 3.5, 0.32, t, size=14.5, bold=True, color=MAROON)
    text(s, M + 3.85, y - 0.02, 4.25, 0.95, b, size=12.5, spacing=1.26)
    y += 1.02
bar(s, M, 6.15, 7.75, 0.018, RULE)
text(s, M, 6.35, 7.75, 0.7,
     "Natural partners span computer vision, wildlife ecology and state agencies; the "
     "corpus is the piece none of them can produce alone.",
     size=12.5, color=MUTE, spacing=1.28)

text(s, 9.35, 1.95, 3.25, 0.3, "Funding directions", size=14.5, bold=True, color=MAROON)
text(s, 9.35, 2.37, 3.25, 3.4,
     ["Candidate targets to scope, not yet selected:",
      "· Federal science agencies supporting AI for environmental monitoring",
      "· Agricultural and natural-resource programmes concerned with deer population "
      "management",
      "· State wildlife agencies that already run transect surveys and carry the "
      "operational need",
      "· Instrumentation and sensor-development programmes for the edge hardware"],
     size=12, spacing=1.26, gap=8)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("saved", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")