#!/usr/bin/env python3
"""Render the deck to PNG without PowerPoint or LibreOffice.

Neither is installed on this cluster, so this walks the .pptx shape tree and redraws
it with PIL. It only has to handle the primitives make_deck.py emits -- solid
rectangles, pictures, and textboxes with per-run size/bold/colour and per-paragraph
alignment, line spacing and space-before.

Calibri is not installed either. DejaVu Sans Condensed is used instead because its
advance widths are close to Calibri's; plain DejaVu Sans is noticeably wider and
would wrap lines the real deck does not.
"""
from __future__ import annotations
import argparse, io, os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

EMU = 914400
FONTS = {
    (False,): "/usr/share/fonts/dejavu-sans-fonts/DejaVuSansCondensed.ttf",
    (True,):  "/usr/share/fonts/dejavu-sans-fonts/DejaVuSansCondensed-Bold.ttf",
}
_cache: dict = {}


def font(pt: float, bold: bool, dpi: int):
    key = (round(pt, 1), bold, dpi)
    if key not in _cache:
        _cache[key] = ImageFont.truetype(FONTS[(bold,)], max(int(pt * dpi / 72), 1))
    return _cache[key]


def wrap(draw, words, f, max_w):
    """Greedy wrap. Explicit newlines are honoured as hard breaks."""
    lines, cur = [], ""
    for token in words.split(" "):
        trial = token if not cur else cur + " " + token
        if draw.textlength(trial, font=f) <= max_w or not cur:
            cur = trial
        else:
            lines.append(cur); cur = token
    lines.append(cur)
    return lines


def render(pptx_path: str, out_dir: str, dpi: int = 140, only=None):
    prs = Presentation(pptx_path)
    SW = prs.slide_width / EMU
    SH = prs.slide_height / EMU
    W, H = int(SW * dpi), int(SH * dpi)
    os.makedirs(out_dir, exist_ok=True)
    written = []

    for idx, slide in enumerate(prs.slides, 1):
        if only and idx not in only:
            continue
        img = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(img)

        for sh in slide.shapes:
            x, y = sh.left / EMU * dpi, sh.top / EMU * dpi
            w, h = sh.width / EMU * dpi, sh.height / EMU * dpi

            if sh.shape_type == 13:                                  # picture
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                img.paste(im.resize((max(int(w), 1), max(int(h), 1)), Image.LANCZOS),
                          (int(x), int(y)))
                continue

            if sh.has_text_frame and sh.text_frame.text.strip():
                tf = sh.text_frame
                cy = y
                for p in tf.paragraphs:
                    runs = [r for r in p.runs if r.text]
                    if not runs:
                        continue
                    pt = max(r.font.size.pt for r in runs if r.font.size) if any(
                        r.font.size for r in runs) else 18
                    bold = bool(runs[0].font.bold)
                    col = runs[0].font.color
                    rgb = tuple(col.rgb) if (col and col.type is not None
                                             and col.rgb is not None) else (0, 0, 0)
                    f = font(pt, bold, dpi)
                    if p.space_before:
                        cy += p.space_before.pt * dpi / 72
                    txt = "".join(r.text for r in runs)
                    lines = []
                    for hard in txt.split("\n"):
                        lines.extend(wrap(d, hard, f, w))
                    lh = pt * 1.22 * (p.line_spacing or 1.0) * dpi / 72
                    for ln in lines:
                        al = str(p.alignment)
                        tw = d.textlength(ln, font=f)
                        lx = x + (w - tw) if "RIGHT" in al else (
                             x + (w - tw) / 2 if "CENTER" in al else x)
                        d.text((lx, cy), ln, font=f, fill=rgb)
                        cy += lh
                continue

            try:                                                     # solid rectangle
                fc = sh.fill.fore_color.rgb
            except Exception:
                continue
            d.rectangle([x, y, x + w, y + h], fill=tuple(fc))

        p = os.path.join(out_dir, f"slide_{idx:02d}.png")
        img.save(p)
        written.append(p)
        print(f"  {os.path.basename(p)}  {W}x{H}")
    return written


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", default="docs/slides/wildlife_results.pptx")
    ap.add_argument("--out", default="docs/slides/png")
    ap.add_argument("--dpi", type=int, default=140)
    ap.add_argument("--slides", default="", help="e.g. 10,11,12,13")
    a = ap.parse_args()
    only = {int(v) for v in a.slides.split(",") if v.strip()} or None
    render(a.pptx, a.out, a.dpi, only)
